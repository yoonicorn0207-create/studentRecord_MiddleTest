
from pptx import Presentation
from pptx.util import Inches, Pt
from pptx.dml.color import RGBColor
from pptx.enum.text import PP_ALIGN

def add_slide(prs, title_text, content_text=None, is_dark=True):
    slide_layout = prs.slide_layouts[1] # Title and Content
    slide = prs.slides.add_slide(slide_layout)
    
    # Set background color to dark
    if is_dark:
        background = slide.background
        fill = background.fill
        fill.solid()
        fill.fore_color.rgb = RGBColor(30, 30, 30)
    
    # Title style
    title = slide.shapes.title
    title.text = title_text
    title_text_frame = title.text_frame
    for paragraph in title_text_frame.paragraphs:
        paragraph.font.color.rgb = RGBColor(255, 255, 255)
        paragraph.font.bold = True
        paragraph.font.size = Pt(32)

    # Content style
    if content_text:
        body_shape = slide.shapes.placeholders[1]
        tf = body_shape.text_frame
        tf.text = content_text
        for paragraph in tf.paragraphs:
            paragraph.font.color.rgb = RGBColor(200, 200, 200)
            paragraph.font.size = Pt(18)
            paragraph.space_after = Pt(10)

def create_presentation():
    prs = Presentation()
    
    # 1. Title Slide
    slide_layout = prs.slide_layouts[0]
    slide = prs.slides.add_slide(slide_layout)
    background = slide.background
    fill = background.fill
    fill.solid()
    fill.fore_color.rgb = RGBColor(20, 20, 20)
    
    title = slide.shapes.title
    title.text = "학생 정보 및 성적 관리 시스템\n개발 프로젝트 결과 보고서"
    for p in title.text_frame.paragraphs:
        p.font.color.rgb = RGBColor(0, 191, 255) # Deep Sky Blue
        p.font.bold = True
        
    subtitle = slide.placeholders[1]
    subtitle.text = "Spring MVC & Docker 기반 성적 관리 솔루션\n발표자: 최창윤"
    for p in subtitle.text_frame.paragraphs:
        p.font.color.rgb = RGBColor(255, 255, 255)

    # Slides List (Total 40)
    content_list = [
        ("목차 (Table of Contents)", "1. 프로젝트 개요\n2. 시스템 아키텍처\n3. 데이터베이스 설계\n4. 백엔드 구현 상세\n5. 핵심 로직: 실시간 석차 계산\n6. 인프라 및 배포 (Docker)\n7. 트러블슈팅 및 해결 경험\n8. 결론 및 향후 계획"),
        
        # Section 1: Overview (3-7)
        ("1.1 프로젝트 배경", "전통적인 성적 관리 방식의 불편함 개선 필요\n교수자 및 관리자가 실시간으로 성적을 입력하고 통계를 확인하는 시스템 구축"),
        ("1.2 프로젝트 목표", "Spring MVC 패턴의 완벽한 이해와 적용\n실시간 석차 산출 로직 구현\nDocker 컨테이너를 활용한 배포 자동화 및 환경 독립성 확보"),
        ("1.3 주요 기능 소개", "학생 성적 CRUD (Create, Read, Update, Delete)\n자동 총점 및 평균 계산\n실시간 전교 석차 자동 갱신 시스템"),
        ("1.4 개발 환경 - Backend", "Framework: Spring MVC 5.x\nLanguage: Java 11\nBuild Tool: Maven\nServer: Apache Tomcat 9.0"),
        ("1.5 개발 환경 - Database & Infra", "DB: MySQL 8.0\nInfra: Docker, Docker-Compose\nFrontend: Node.js (Static Server)"),
        
        # Section 2: Architecture (8-12)
        ("2.1 시스템 아키텍처 개요", "유지보수와 확장성을 고려한 계층형 아키텍처 적용\nSpring MVC 패턴을 통한 관심사 분리(SoC)"),
        ("2.2 Spring MVC 흐름도", "DispatcherServlet -> HandlerMapping -> Controller -> Service -> DAO/JDBC -> ViewResolver"),
        ("2.3 디렉토리 구조 분석", "src/main/java: 비즈니스 로직 및 컨트롤러\nsrc/main/webapp: JSP 및 설정 파일\nfrontend-node: 프론트엔드 서버 환경"),
        ("2.4 데이터 흐름 (Data Flow)", "Client Request -> JSON Payload -> DTO Mapping -> Service Processing -> DB Persistence -> JSON/JSP Response"),
        ("2.5 프론트엔드-백엔드 연동", "REST API 방식의 통신\nNode.js 서버와 Spring 서버 간의 비동기 AJAX 통신 구조"),
        
        # Section 3: Database (13-17)
        ("3.1 데이터베이스 모델링", "Table Name: student_scores\n학생 식별자(ID)를 중심으로 성적 및 통계 데이터를 단일 테이블로 관리"),
        ("3.2 스키마 정의 (DDL)", "id (INT, PK, AI)\nstudent_name (VARCHAR)\nkor, eng, math (INT)\ntotal_score, avg_score, student_rank"),
        ("3.3 데이터 영속성 전략", "JDBC PreparedStatement를 통한 보안성 강화 (SQL Injection 방지)\nDataSource 커넥션 풀 활용으로 성능 최적화"),
        ("3.4 init.sql 자동화", "Docker 컨테이너 실행 시 스크립트를 통한 테이블 자동 생성\n배포 즉시 서비스 운영이 가능한 'Infrastructure as Code' 실현"),
        ("3.5 데이터 정렬 및 조회 쿼리", "ORDER BY student_rank ASC, id DESC\n석차순 정렬을 기본으로 하되 동일 석차 시 최신순 조회"),
        
        # Section 4: Implementation (18-27)
        ("4.1 GradeDTO (Model)", "데이터 전송 객체 구현\nGetter/Setter를 통한 캡슐화 준수\nJSON 직렬화를 위한 필드 구성"),
        ("4.2 recordcardController (Controller)", "@Controller, @RequestMapping 활용\nRESTful API 설계를 통한 확장성 확보"),
        ("4.3 @CrossOrigin 적용", "프론트엔드와 백엔드 포트 불일치 해결\n브라우저 보안 정책(CORS)에 대응하는 설정 추가"),
        ("4.4 API: 학생 추가 (/add)", "@RequestBody를 활용한 JSON 데이터 수신\nService 계층으로 비즈니스 로직 위임"),
        ("4.5 API: 학생 목록 조회 (/list)", "Service로부터 List<GradeDTO> 수신\n@ResponseBody를 통한 JSON 배열 반환"),
        ("4.6 GradeService (Service Layer)", "@Service 어노테이션을 통한 컴포넌트 스캔 대상 등록\n비즈니스 로직의 중심부"),
        ("4.7 트랜잭션 및 JDBC 연동", "Connection, PreparedStatement의 안전한 자원 반납(Try-with-resources)\n예외 처리 로직 구현"),
        ("4.8 성적 계산 로직 상세", "입력된 국/영/수 점수 합산\n소수점 둘째 자리까지의 평균 자동 산출 (Math.round 활용)"),
        ("4.9 API: 학생 수정 및 삭제", "ID 기반의 데이터 갱신 및 삭제 처리\n작업 완료 후 반드시 석차 재계산 로직 호출"),
        ("4.10 Spring Bean 설정 (XML)", "beans_product.xml을 통한 DataSource 및 Property 관리\nComponent Scan 설정"),
        
        # Section 5: Rank Logic (28-30)
        ("5.1 실시간 석차 계산 로직 (UpdateAllRanks)", "데이터 변경 시마다 전체 학생 등수를 전수 조사\n점수 내림차순 정렬 후 순차적 갱신"),
        ("5.2 동점자 처리 알고리즘", "이전 학생 점수와 비교하여 같을 경우 동일 등수 부여\n다를 경우 현재 처리 건수(count)를 등수로 설정"),
        ("5.3 석차 업데이트 성능 고려", "단일 트랜잭션 내 처리\n데이터 규모 확장을 대비한 인덱싱 고려 필요성 확인"),
        
        # Section 6: Infrastructure (31-35)
        ("6.1 Docker 컨테이너 전략", "애플리케이션 개발 환경과 운영 환경의 일치(Consistency)\n컨테이너 기반 가상화의 장점 활용"),
        ("6.2 Backend Dockerfile 분석", "Tomcat 9.0 베이스 이미지 활용\nMaven 빌드 결과물(war) 자동 복사 및 배포"),
        ("6.3 Docker-Compose 오케스트레이션", "DB, BE, FE 서비스를 하나의 네트워크(student-network)로 통합\ndepends_on을 이용한 실행 순서 제어"),
        ("6.4 데이터 볼륨 매핑 (Persistence)", "db_data 볼륨을 통해 컨테이너 삭제 후에도 DB 데이터 보존\n로컬 디렉토리와의 동기화 설정"),
        ("6.5 네트워크 인프라 설계", "내부 브릿지 네트워크를 통한 서비스 간 격리 및 보안 강화\n포트 포워딩(8080, 3000, 3306) 설정"),
        
        # Section 7: Troubleshooting (36-39)
        ("7.1 트러블슈팅: CORS 에러 발생", "문제: 포트가 다른 프론트엔드에서 API 요청 시 브라우저가 차단\n해결: Controller에 @CrossOrigin 어노테이션 추가 및 특정 도메인 허용"),
        ("7.2 트러블슈팅: DB 접속 타이밍 이슈", "문제: Docker-Compose 실행 시 DB가 미처 준비되기 전 BE가 실행되어 에러 발생\n해결: depends_on 설정 및 애플리케이션 내 재시도 로직(Retry) 검토"),
        ("7.3 트러블슈팅: 동점자 석차 오류", "문제: 1등이 2명일 때 다음 학생이 2등으로 표기되는 문제\n해결: rank 변수와 count 변수를 분리하여 표준 석차 방식(1, 1, 3등)으로 알고리즘 수정"),
        ("7.4 트러블슈팅: 한글 깨짐 현상", "문제: DB 저장 시 학생 이름이 ???로 표시됨\n해결: JDBC URL에 characterEncoding=UTF-8 옵션 추가 및 MySQL 캐릭터셋 설정 변경"),
        
        # Section 8: Conclusion (40)
        ("8.1 결론 및 배운 점", "Spring MVC의 전반적인 라이프사이클 체득\nDocker를 통한 DevOps 프로세스의 중요성 인식\n비즈니스 로직 구현 시 예외 상황 처리의 세밀함 배양\n감사합니다. Q&A"),
    ]

    for title, content in content_list:
        add_slide(prs, title, content)

    # Save to Desktop
    import os
    desktop_path = os.path.join(os.path.join(os.environ['USERPROFILE']), 'Desktop')
    save_path = os.path.join(desktop_path, "학생성적관리시스템_발표자료.pptx")
    prs.save(save_path)
    print(f"Presentation saved to: {save_path}")

if __name__ == "__main__":
    create_presentation()
